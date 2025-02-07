import pandas as pd
import json
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders.csv_loader import CSVLoader
from pypdf import PdfReader
from langchain_core.documents import Document
import logging
import os
import asyncio


async def load_pdf(file_path):
    try:
        logging.info(f"Starting to load PDF: {file_path}")

        # Verify file exists and is readable
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")

        def read_pdf():
            reader = PdfReader(file_path)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():  # Only include pages with content
                    pages.append(
                        Document(
                            page_content=text,
                            metadata={"source": file_path, "page": i}
                        )
                    )
            return pages

        # Run PDF reading in a thread pool to avoid blocking
        pages = await asyncio.get_event_loop().run_in_executor(None, read_pdf)

        if not pages:
            logging.error(f"No valid pages found in {file_path}")
            return None

        logging.info(
            f"Successfully loaded {len(pages)} pages from {file_path}")
        logging.info(f"First page metadata: {pages[0].metadata}")
        logging.info(
            f"First page content sample: {pages[0].page_content[:200]}...")

        return pages
    except Exception as e:
        logging.error(
            f"Error loading PDF {file_path}: {str(e)}", exc_info=True)
        return None


async def load_py(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            content = f.read()
            return content.strip()
    except Exception as e:
        print(f"Error loading PY: {str(e)}")
        return None


async def load_docx(file):
    try:
        loader = Docx2txtLoader(file)
        data = loader.load()
        print(data)
        return data[0].page_content
    except Exception as e:
        print(f"Error loading DOCX: {str(e)}")
        return None


async def load_txt(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except Exception as e:
        print(f"Error loading TXT: {str(e)}")
        return None


async def load_md(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            md_text = f.read()
            html = markdown.markdown(md_text)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text().strip()
    except Exception as e:
        print(f"Error loading MD: {str(e)}")
        return None


async def load_html(file_path: str) -> str:
    """Load and process HTML file content"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Parse HTML with BeautifulSoup
        soup = BeautifulSoup(content, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text content
        text = soup.get_text()

        # Break into lines and remove leading/trailing space
        lines = (line.strip() for line in text.splitlines())

        # Break multi-headlines into a line each
        chunks = (phrase.strip()
                  for line in lines for phrase in line.split("  "))

        # Drop blank lines
        text = ' '.join(chunk for chunk in chunks if chunk)

        return text
    except Exception as e:
        logging.error(f"Error loading HTML file {file_path}: {str(e)}")
        return None


async def load_csv(file):
    try:
        loader = CSVLoader(file)
        data = loader.load()
        return data
    except Exception as e:
        print(f"Error loading CSV: {str(e)}")
        return None


async def load_json(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2)
    except Exception as e:
        print(f"Error loading JSON: {str(e)}")
        return None


def load_pptx(file):
    try:
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
    except Exception as e:
        print(f"Error loading PPTX: {str(e)}")
        return None


def load_xlsx(file):
    try:
        df = pd.read_excel(file)
        return df.to_string().strip()
    except Exception as e:
        print(f"Error loading XLSX: {str(e)}")
        return None


async def load_docx(file):
    try:
        # Run the synchronous loader in a thread pool to avoid blocking
        def load_docx_sync():
            loader = Docx2txtLoader(file)
            data = loader.load()
            return data[0].page_content if data else None

        content = await asyncio.get_event_loop().run_in_executor(None, load_docx_sync)
        if content:
            logging.info(f"Successfully loaded DOCX file: {file}")
            return content
        return None
    except Exception as e:
        logging.error(f"Error loading DOCX: {str(e)}")
        return None
